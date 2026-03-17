# -*- coding: utf-8 -*-

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DESKTOP = Path(r"D:/Desktop") if Path(r"D:/Desktop").exists() else Path.home() / "Desktop"
SOURCE_DOCX = next(DESKTOP.glob("115-20232115212098*.docx"))
PPT_PATH = DESKTOP / "婴儿癫痫性痉挛综合征答辩PPT-胡梦碟.pptx"
SCRIPT_DOCX_PATH = DESKTOP / "婴儿癫痫性痉挛综合征答辩演讲稿-胡梦碟.docx"
SCRIPT_TXT_PATH = DESKTOP / "婴儿癫痫性痉挛综合征答辩演讲稿-胡梦碟.txt"

TITLE = "婴儿癫痫性痉挛综合征临床特点及不同药物治疗方案的疗效分析"
AUTHOR = "胡梦碟"
ADVISOR = "吴丽文 教授"
MAJOR = "儿科学（小儿神经）"
SCHOOL = "南华大学儿科学院"
TIME = "2026年4月"
HOSPITAL = "湖南省儿童医院神经内科"

NAVY = RGBColor(17, 46, 81)
TEAL = RGBColor(0, 150, 149)
CYAN = RGBColor(74, 181, 195)
ORANGE = RGBColor(234, 131, 62)
RED = RGBColor(204, 74, 62)
SLATE = RGBColor(70, 85, 103)
MUTED = RGBColor(104, 116, 132)
LIGHT_BG = RGBColor(245, 248, 252)
CARD_BG = RGBColor(255, 255, 255)
LINE = RGBColor(218, 226, 235)
SOFT_BLUE = RGBColor(228, 240, 249)

SLIDE_SCRIPTS = []


def pct_reduction(base_value, new_value):
    return round((base_value - new_value) / base_value * 100, 1)


def set_font(run, size, color, bold=False, name="Microsoft YaHei"):
    font = run.font
    font.name = name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color


def fill_shape(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=NAVY,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.08,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_font(run, size, color, bold=bold)
    return box


def add_bullets(slide, left, top, width, height, items, size=14, text_color=SLATE, accent=TEAL):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(8)
        prefix = paragraph.add_run()
        prefix.text = "• "
        set_font(prefix, size, accent, bold=True)
        content = paragraph.add_run()
        content.text = item
        set_font(content, size, text_color)

    return box


def add_header(slide, section, title, subtitle=None):
    fill_shape(
        slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.76)
        ),
        NAVY,
    )
    fill_shape(
        slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(0.88), Inches(1.05), Inches(0.08)
        ),
        TEAL,
    )
    add_textbox(slide, 0.55, 0.14, 2.5, 0.28, section, size=11, color=RGBColor(208, 225, 243), bold=True)
    add_textbox(slide, 0.55, 0.98, 8.3, 0.5, title, size=25, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, 0.55, 1.46, 8.8, 0.3, subtitle, size=10.5, color=MUTED)
    fill_shape(
        slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(12.98), Inches(0.76), Inches(0.08), Inches(6.35)
        ),
        SOFT_BLUE,
    )


def add_page_no(slide, index):
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.85), Inches(0.16), Inches(0.85), Inches(0.28)
    )
    fill_shape(chip, TEAL)
    add_textbox(
        slide,
        11.89,
        0.18,
        0.77,
        0.2,
        f"{index:02d}",
        size=10,
        color=RGBColor(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.01,
    )


def add_card(slide, left, top, width, height, title, body_items, accent=TEAL):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = LINE
    fill_shape(
        slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left + 0.18),
            Inches(top + 0.18),
            Inches(0.95),
            Inches(0.24),
        ),
        accent,
    )
    add_textbox(slide, left + 0.18, top + 0.52, width - 0.36, 0.35, title, size=14, color=NAVY, bold=True)
    add_bullets(slide, left + 0.16, top + 0.9, width - 0.32, height - 1.02, body_items, size=13, accent=accent)


def metric_value_size(value, width):
    compact = str(value).replace(" ", "").replace("\n", "")
    length = len(compact)
    if width <= 2.0:
        return 16
    if length >= 18:
        return 13.5
    if length >= 14:
        return 15
    if length >= 10 or width <= 2.4:
        return 18
    if width <= 3.0:
        return 20
    return 24


def add_metric_card(slide, left, top, width, height, caption, value, note, accent=TEAL, value_color=NAVY):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = LINE
    fill_shape(
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height)),
        accent,
    )
    caption_size = 10.5 if width < 2.4 else 11.5
    note_size = 8.6 if width < 3.0 else 10
    note_height = 0.28 if height < 1.45 else 0.34
    note_top = top + height - note_height - 0.08
    value_top = top + 0.46
    value_height = max(0.34, note_top - value_top - 0.05)
    add_textbox(
        slide,
        left + 0.24,
        top + 0.16,
        width - 0.36,
        0.24,
        caption,
        size=caption_size,
        color=MUTED,
        bold=True,
        margin=0.02,
    )
    add_textbox(
        slide,
        left + 0.24,
        value_top,
        width - 0.36,
        value_height,
        value,
        size=metric_value_size(value, width),
        color=value_color,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.02,
    )
    add_textbox(
        slide,
        left + 0.24,
        note_top,
        width - 0.36,
        note_height,
        note,
        size=note_size,
        color=SLATE,
        margin=0.02,
    )


def add_data_note(slide, left, top, text):
    add_textbox(slide, left, top, 5.7, 0.25, f"数据来源：{text}", size=9.2, color=MUTED)


def record_script(title, points, script):
    SLIDE_SCRIPTS.append({"title": title, "points": points, "script": script})


def add_table(slide, left, top, width, height, rows, col_widths=None, font_size=10.5):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    if col_widths:
        for idx, col_width in enumerate(col_widths):
            table.columns[idx].width = Inches(col_width)

    for row_index, row in enumerate(rows):
        for col_index, cell_value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(cell_value)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            set_font(run, font_size, NAVY if row_index == 0 else SLATE, bold=row_index == 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT_BLUE if row_index == 0 else CARD_BG

    return table


def add_column_chart(slide, left, top, width, height, categories, series_data, max_value=100):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_data:
        chart_data.add_series(series_name, values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = max_value
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.value_axis.format.line.color.rgb = LINE
    chart.category_axis.format.line.color.rgb = LINE
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.font.size = Pt(9)

    palette = [NAVY, TEAL, ORANGE, CYAN, RED]
    for index, series in enumerate(chart.series):
        color = palette[index % len(palette)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.color.rgb = color


def add_bar_chart(slide, left, top, width, height, categories, values, title=None, color=TEAL, max_value=None):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("数值", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    ).chart
    chart.has_legend = False
    if max_value is not None:
        chart.value_axis.maximum_scale = max_value
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.value_axis.format.line.color.rgb = LINE
    chart.category_axis.format.line.color.rgb = LINE
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.font.size = Pt(9)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = color
    chart.series[0].format.line.color.rgb = color
    if title:
        add_textbox(slide, left, top - 0.28, width, 0.22, title, size=11.5, color=NAVY, bold=True)


def add_timeline_item(slide, left, top, label, description, accent):
    fill_shape(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(0.18), Inches(0.18)), accent)
    add_textbox(slide, left + 0.26, top - 0.03, 2.4, 0.24, label, size=12.2, color=NAVY, bold=True)
    add_textbox(slide, left + 0.26, top + 0.18, 2.65, 0.58, description, size=10.8, color=SLATE)


def add_section_label(slide, left, top, text, color=TEAL):
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(1.25), Inches(0.28)
    )
    fill_shape(label, color)
    add_textbox(
        slide,
        left + 0.02,
        top + 0.02,
        1.21,
        0.2,
        text,
        size=9.5,
        color=RGBColor(255, 255, 255),
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.01,
    )


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = TITLE
    prs.core_properties.author = AUTHOR
    prs.core_properties.subject = "硕士论文答辩PPT"
    prs.core_properties.comments = f"根据论文《{TITLE}》自动生成"
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    fill_shape(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.18)), NAVY)
    fill_shape(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.92), Inches(0), Inches(4.413), Inches(7.5)), RGBColor(232, 245, 246))
    add_textbox(slide, 0.72, 1.38, 7.55, 1.45, TITLE, size=26, color=NAVY, bold=True)
    add_textbox(slide, 0.74, 2.95, 5.2, 0.3, "硕士学位论文答辩汇报", size=13.5, color=TEAL, bold=True)
    add_textbox(slide, 0.74, 3.45, 3.0, 0.25, f"汇报人：{AUTHOR}", size=12.4, color=SLATE)
    add_textbox(slide, 0.74, 3.82, 3.4, 0.25, f"导师：{ADVISOR}", size=12.4, color=SLATE)
    add_textbox(slide, 0.74, 4.19, 4.2, 0.25, f"专业：{MAJOR}", size=12.4, color=SLATE)
    add_textbox(slide, 0.74, 4.56, 4.2, 0.25, f"学院：{SCHOOL}", size=12.4, color=SLATE)
    add_textbox(slide, 0.74, 4.93, 2.0, 0.25, TIME, size=12.4, color=SLATE)
    add_metric_card(slide, 9.15, 1.5, 3.45, 1.3, "研究样本", "226例 IESS患儿", "2020.01 - 2025.12，单中心回顾性队列")
    add_metric_card(slide, 9.15, 3.0, 3.45, 1.3, "最关键结果", "85.5% vs 70.6%", "泼尼松和/或氨己烯酸组总有效率更高", accent=ORANGE)
    add_metric_card(slide, 9.15, 4.5, 3.45, 1.3, "复发比较", "27.4% vs 50.9%", "6个月复发率明显更低", accent=CYAN)
    record_script("封面", ["说明论文题目、作者、导师和答辩时间", "提前引出样本量、有效率和复发率三项核心信息"], "各位老师好，我今天汇报的题目是《婴儿癫痫性痉挛综合征临床特点及不同药物治疗方案的疗效分析》。本研究纳入226例IESS患儿，比较ACTH与泼尼松联合或不联合氨己烯酸两类方案的疗效、复发、神经发育及经济学差异。先给出最核心的结果，泼尼松和或氨己烯酸组总有效率更高，6个月复发率更低，同时住院时间和费用也更少。")

    # 2. Outline
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "汇报目录", "本次答辩共分五个部分", "按研究逻辑依次展开：背景、方法、结果、讨论、结论")
    add_card(slide, 0.78, 2.0, 2.4, 3.85, "一、研究背景", ["IESS疾病特点", "一线治疗现状", "临床争议与研究意义"], accent=TEAL)
    add_card(slide, 3.42, 2.0, 2.4, 3.85, "二、研究方法", ["研究对象与分组", "治疗方案", "观察指标与统计学"], accent=ORANGE)
    add_card(slide, 6.06, 2.0, 2.4, 3.85, "三、研究结果", ["短期疗效", "复发、安全性与费用", "发育结局与影响因素"], accent=CYAN)
    add_card(slide, 8.7, 2.0, 2.4, 3.85, "四、讨论", ["结果解释", "与既往研究对照", "临床启示"], accent=RED)
    add_card(slide, 11.34, 2.0, 1.22, 3.85, "五、结论", ["总结", "局限性", "展望"], accent=TEAL)
    add_page_no(slide, 2)
    record_script("汇报目录", ["交代汇报结构", "提示老师后续重点会放在疗效、复发和治疗时机"], "汇报主要分为五部分。第一部分是研究背景，说明为什么要重新比较IESS的一线治疗方案。第二部分是研究方法，包括对象来源、分组、治疗方案和观察指标。第三部分是研究结果，这一部分是今天的重点。第四部分对结果做讨论和临床解释。最后汇报研究结论、局限性和后续展望。")

    # 3. Background
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究背景", "IESS 是婴幼儿期严重癫痫性脑病", "疾病特点决定了治疗不能只看止痉，还必须重视发育与长期预后")
    add_card(slide, 0.72, 1.98, 4.0, 4.9, "疾病概述", ["IESS即婴儿癫痫性痉挛综合征，2022年ILAE继续沿用该综合征命名。", "起病年龄多在1至24月龄，发病高峰为3至12月龄。", "典型表现包括成串痉挛、高度失律脑电图以及不同程度发育停滞或倒退。"])
    add_card(slide, 4.98, 1.98, 3.8, 4.9, "疾病负担", ["IESS虽属相对少见疾病，但对神经发育影响显著。", "持续异常放电可损害发育中的大脑，延迟治疗可能带来不可逆后果。", "远期仍可残留认知、运动障碍，甚至演变为其他癫痫综合征。"], accent=ORANGE)
    add_card(slide, 9.02, 1.98, 3.45, 4.9, "治疗目标", ["尽快实现痉挛停止。", "同步改善脑电图高度失律，达到电-临床缓解。", "尽量降低复发风险，并保护神经发育轨迹。"], accent=CYAN)
    add_page_no(slide, 3)
    record_script("研究背景：疾病概况", ["说明IESS定义、起病年龄和典型临床表现", "强调其对神经发育的长期损害"], "IESS是婴幼儿期严重的癫痫性脑病，既往常被称为婴儿痉挛症或West综合征。它的典型特点包括痉挛发作、脑电图高度失律以及不同程度的神经发育受损，起病高峰主要集中在3到12月龄。IESS的临床管理之所以特殊，是因为它不仅是一种癫痫发作问题，更会深刻影响患儿后续的认知、语言和运动发育。因此，治疗目标不能只停留在止痉，还必须把发育保护放到同等重要的位置。")

    # 4. Current treatment
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究背景", "当前一线治疗方案已明确，但最佳临床选择仍存在争议", "争议焦点集中在疗效差异、复发控制、可及性与经济负担")
    add_card(slide, 0.72, 1.98, 4.0, 4.9, "共识部分", ["国际与国内指南普遍将激素和氨己烯酸视为IESS一线治疗。", "ACTH、口服糖皮质激素以及联合VGB均有较多循证支持。", "普遍认为尽早治疗、约2周评估反应是关键策略。"])
    add_card(slide, 4.98, 1.98, 4.0, 4.9, "ACTH方案特点", ["优点：传统证据充分，在部分研究中对脑电图改善更有优势。", "问题：通常需住院静脉给药，流程复杂，成本高，不良反应相对较多。", "现实中可能受到费用和可及性限制。"], accent=ORANGE)
    add_card(slide, 9.24, 1.98, 3.25, 4.9, "口服激素 ± VGB", ["优点：给药便捷，启动更快，临床更易推广。", "部分研究提示短期止痉和联合治疗效果并不逊于ACTH。", "仍需真实世界数据验证疗效、复发与发育获益。"], accent=CYAN)
    add_page_no(slide, 4)
    record_script("研究背景：治疗现状与争议", ["ACTH与口服激素±VGB都属于一线治疗", "真正的临床争议在于综合效益谁更优"], "从目前循证证据看，ACTH、口服糖皮质激素以及氨己烯酸都属于IESS的一线治疗工具。ACTH的传统证据较充分，但它往往需要住院静脉滴注，治疗成本更高，不良反应也相对更多。相较之下，泼尼松联合或不联合氨己烯酸的方案在临床上更便捷、更容易启动，也更符合实际可及性需求。因此，真正的争议不是两者是否都有效，而是哪一种方案在疗效、复发、发育和成本这些综合维度上更值得优先采用。")

    # 5. Purpose
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究目的", "本研究要回答的核心问题", "基于真实世界临床资料，比较不同治疗路径的综合价值")
    add_metric_card(slide, 0.86, 2.0, 3.75, 1.55, "问题一", "哪种方案短期疗效更好？", "观察2周无痉挛、电-临床缓解与总体有效率")
    add_metric_card(slide, 4.8, 2.0, 3.75, 1.55, "问题二", "哪种方案更能降低复发？", "随访初始缓解患儿6个月内是否再次发作", accent=ORANGE)
    add_metric_card(slide, 8.74, 2.0, 3.75, 1.55, "问题三", "药物选择是否影响预后？", "结合发育、治疗时机和回归分析综合判断", accent=CYAN)
    add_card(slide, 1.12, 4.15, 11.1, 2.05, "研究价值", ["如果口服方案在疗效和复发控制上不劣于甚至优于ACTH，同时又能减少住院和费用，那么它将对IESS临床路径优化具有直接意义。"], accent=TEAL)
    add_page_no(slide, 5)
    record_script("研究目的", ["围绕疗效、复发和预后三个维度提出问题", "强调这是临床决策导向的研究"], "本研究主要聚焦三个问题。第一，ACTH和泼尼松和或氨己烯酸相比，哪一种在短期内更容易控制痉挛发作。第二，在初始缓解后，哪一种方案复发更少。第三，药物选择、发作年龄和治疗时机这些因素中，哪些因素真正与疗效和发育结局更相关。换句话说，这项研究的目标并不是单纯比较一种药物，而是希望为一线临床决策提供更实用的依据。")

    # 6. Design
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究方法", "研究设计与病例来源", "单中心回顾性研究，所有病例均来自湖南省儿童医院")
    add_metric_card(slide, 0.85, 2.0, 2.7, 1.55, "研究类型", "回顾性队列研究", "单中心真实世界临床资料")
    add_metric_card(slide, 3.78, 2.0, 2.7, 1.55, "研究时间", "2020.01 - 2025.12", "随访截止至2026.01.30", accent=ORANGE)
    add_metric_card(slide, 6.71, 2.0, 2.7, 1.55, "研究地点", HOSPITAL, "病例资料来自住院、门诊和电话随访", accent=CYAN)
    add_metric_card(slide, 9.64, 2.0, 2.7, 1.55, "伦理审批", "HCHLL-2026-110", "研究已通过医学伦理审批", accent=RED)
    add_card(slide, 0.9, 4.0, 11.9, 2.15, "技术路线", ["病例筛选 -> 分组 -> 治疗后2周评估短期疗效 -> 随访6个月评估复发与发育 -> 结合1年发育结局与Logistic回归分析影响因素"], accent=TEAL)
    add_page_no(slide, 6)
    record_script("研究设计与病例来源", ["介绍研究时间、地点和设计类型", "说明病例来自住院和随访记录"], "本研究为单中心回顾性队列研究，病例均来自湖南省儿童医院神经内科。研究时间跨度为2020年1月至2025年12月，随访截止到2026年1月30日。病例的临床资料主要来自住院病历，复发和长期结局则结合门诊随访及电话回访获取。研究已经通过医院伦理审批，因此具备规范的研究基础。")

    # 7. Criteria
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究方法", "研究对象、纳入标准与分组方式", "最终纳入226例患儿，按首选治疗方案分为两组")
    add_card(slide, 0.72, 1.95, 3.9, 4.85, "纳入标准", ["符合ILAE婴儿癫痫性痉挛综合征诊断标准。", "首次痉挛发作年龄小于2岁。", "随访时间大于3个月。"], accent=TEAL)
    add_card(slide, 4.84, 1.95, 3.9, 4.85, "排除标准", ["合并重要脏器功能不全。", "合并免疫缺陷疾病。", "存在激素或氨己烯酸用药禁忌证。"], accent=ORANGE)
    add_card(slide, 8.96, 1.95, 3.4, 4.85, "分组情况", ["ACTH组：148例。", "泼尼松和/或氨己烯酸组：78例。", "按初始首选治疗方案进行分组。"], accent=CYAN)
    add_page_no(slide, 7)
    record_script("研究对象与纳排标准", ["说明病例为什么能进入研究", "突出最终是按首选治疗方案分组"], "研究对象需要同时满足三个条件，也就是符合IESS诊断标准、首次痉挛发作年龄小于2岁，并且至少有3个月以上随访。对于存在重要脏器功能不全、免疫缺陷，或激素、氨己烯酸禁忌证的患儿予以排除。最终符合标准的共有226例，其中ACTH组148例，泼尼松和或氨己烯酸组78例。这里的分组依据是初始首选治疗方案，而不是后续调整方案。")

    # 8. Regimens
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究方法", "两组治疗方案与随访节点", "围绕2周、6个月和1年三个时间点评估疗效与预后")
    add_card(slide, 0.72, 1.95, 4.15, 4.3, "ACTH方案", ["合成ACTH 20 - 40 IU/天静脉滴注。", "初始治疗2周，每日1次。", "若2周内痉挛完全控制，可序贯口服泼尼松维持。"], accent=TEAL)
    add_card(slide, 5.08, 1.95, 4.15, 4.3, "泼尼松和/或氨己烯酸方案", ["泼尼松 8 mg/kg/天，分3次，最大60 mg/天。", "氨己烯酸 50 - 100 mg/kg/天，必要时可增至150 mg/kg/天。", "控制后逐渐减量，每周减少25%。"], accent=ORANGE)
    add_section_label(slide, 9.65, 2.0, "随访节点", color=CYAN)
    fill_shape(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.8), Inches(3.18), Inches(2.2), Inches(0.05)), LINE)
    add_timeline_item(slide, 9.78, 2.85, "治疗后2周", "评估痉挛是否停止、电-临床是否缓解", TEAL)
    add_timeline_item(slide, 10.52, 4.1, "治疗后6个月", "评估复发、不良反应及神经发育分级", ORANGE)
    add_timeline_item(slide, 11.28, 5.35, "治疗后1年", "继续追踪神经发育结局", CYAN)
    add_page_no(slide, 8)
    record_script("治疗方案与随访节点", ["简要说明两组用药路径", "突出研究的关键随访时间点"], "在治疗方案上，ACTH组主要采取低剂量静脉滴注的ACTH方案，疗程2周，若痉挛完全控制再序贯口服泼尼松维持。而泼尼松和或氨己烯酸组则以口服高剂量泼尼松为基础，必要时联合氨己烯酸，并在控制后逐渐减量。疗效评估主要围绕三个时间点展开：治疗后2周看短期止痉和电-临床缓解，6个月看复发和发育，1年继续追踪发育结局。")

    # 9. Outcomes
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究方法", "观察指标与统计学分析", "把短期疗效、复发、发育与安全性纳入同一评价框架")
    add_card(slide, 0.72, 1.95, 3.7, 4.85, "主要疗效指标", ["2周内短期完全缓解率。", "2周内痉挛发作控制率。", "总体疗效分为无效、有效、显效。"], accent=TEAL)
    add_card(slide, 4.64, 1.95, 3.75, 4.85, "次要疗效指标", ["2周EEG缓解率。", "6个月复发率。", "6个月及1年神经发育分级。"], accent=ORANGE)
    add_card(slide, 8.61, 1.95, 3.75, 4.85, "统计学方法", ["计量资料用均数±标准差表示。", "计数资料采用χ²检验或Fisher精确检验。", "采用二元Logistic回归分析疗效独立影响因素。"], accent=CYAN)
    add_page_no(slide, 9)
    record_script("观察指标与统计学分析", ["解释研究到底比较了哪些终点", "说明为什么还需要Logistic回归"], "本研究将疗效评价分成三个层面。第一是短期层面，包括2周内是否达到痉挛停止、是否达到电-临床完全缓解，以及总体疗效分级。第二是中期层面，重点看6个月内是否复发。第三是预后层面，关注6个月和1年的神经发育分级。同时，也记录治疗期间的不良反应和住院经济学指标。统计学上先做组间比较，再用Logistic回归分析哪些变量是疗效的独立影响因素。")

    # 10. Baseline and etiology
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "基线资料总体可比，病因分布以病因不明、结构性和遗传性为主", "为后续疗效比较提供了较好的组间可比性")
    add_table(slide, 0.75, 2.0, 5.8, 3.25, [["指标", "ACTH组", "泼尼松和/或氨己烯酸组", "P值"], ["样本量", "148", "78", "-"], ["年龄（月）", "6.92 ± 4.912", "6.24 ± 5.245", "0.807"], ["<12月发作", "130", "68", "0.886"], ["<28天启动治疗", "92", "48", "0.927"], ["总不良反应", "40.5%", "28.2%", "0.066"]], col_widths=[1.3, 1.35, 2.2, 0.95], font_size=10)
    add_bar_chart(slide, 6.95, 2.35, 5.05, 3.55, ["病因不明", "结构性", "遗传性", "感染性", "代谢性"], [115, 65, 40, 3, 3], title="病因构成（例）", color=CYAN, max_value=130)
    add_data_note(slide, 0.78, 6.15, "表3.1、图3.1")
    add_page_no(slide, 10)
    record_script("基线特征与病因构成", ["两组大部分基线指标无统计学差异", "病因分布以病因不明、结构性、遗传性为主"], "在基线资料比较中，两组在年龄、发作年龄、起病到治疗启动时间等大部分指标上均无统计学差异，提示两组总体具有较好的可比性。病因方面，以病因不明最多，占50.9%；其次是结构性，占28.8%；再次是遗传性，占17.7%。这一分布与既往研究基本一致，也提示IESS在病因学上具有明显异质性。")

    # 11. EEG remission
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "电-临床缓解率两组相近", "脑电图高度失律缓解并未出现显著组间差异")
    add_metric_card(slide, 0.9, 2.0, 2.6, 1.5, "纳入分析", "135例", "仅脑电图表现为高度失律者纳入该分析")
    add_metric_card(slide, 3.8, 2.0, 2.8, 1.5, "ACTH组", "66.0%", "完全缓解率", accent=TEAL)
    add_metric_card(slide, 6.95, 2.0, 2.8, 1.5, "泼尼松和/或氨己烯酸组", "71.4%", "完全缓解率", accent=ORANGE)
    add_metric_card(slide, 10.1, 2.0, 2.25, 1.5, "统计学", "P = 0.555", "差异无统计学意义", accent=CYAN)
    add_card(slide, 1.0, 4.1, 11.3, 2.0, "结果解读", ["虽然泼尼松和/或氨己烯酸组的电-临床缓解率略高，但两组差异不显著。说明脑电图改善并未完全跟随短期止痉优势同步放大。"], accent=TEAL)
    add_data_note(slide, 0.98, 6.35, "表3.2、图3.2")
    add_page_no(slide, 11)
    record_script("电-临床缓解", ["仅对脑电图高度失律患儿进行分析", "两组电-临床缓解率差异不显著"], "对于脑电图表现为高度失律的135例患儿，我们比较了治疗2周后的电-临床完全缓解情况。结果显示，ACTH组完全缓解率为66.0%，泼尼松和或氨己烯酸组为71.4%，两组差异没有统计学意义。也就是说，口服方案在短期止痉上更有优势，但这一优势在脑电图完全缓解这个指标上并没有被显著放大。")

    # 12. Short-term control and efficacy
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "短期痉挛控制与总体疗效：口服方案更优", "2周无痉挛率和总体有效率均优于ACTH组")
    add_column_chart(slide, 0.75, 2.1, 7.0, 4.2, ["2周无痉挛率", "总体有效率"], [("ACTH组", [61.6, 70.6]), ("泼尼松和/或氨己烯酸组", [84.7, 85.5])], max_value=100)
    add_metric_card(slide, 8.25, 2.0, 4.0, 1.3, "2周无痉挛率", "84.7% vs 61.6%", "χ² = 12.06，P = 0.005", accent=TEAL)
    add_metric_card(slide, 8.25, 3.65, 4.0, 1.3, "总体有效率", "85.5% vs 70.6%", "χ² = 6.102，P = 0.047", accent=ORANGE)
    add_metric_card(slide, 8.25, 5.3, 4.0, 1.3, "临床含义", "快速止痉更占优", "提示口服方案在短期疗效方面具有现实优势", accent=CYAN)
    add_data_note(slide, 0.82, 6.45, "表3.3、表3.8、图3.3")
    add_page_no(slide, 12)
    record_script("短期痉挛控制与总体疗效", ["2周无痉挛率显著更高", "总体有效率也更高", "这是本研究最直接的疗效证据"], "这是本研究最关键的结果之一。2周随访时，泼尼松和或氨己烯酸组的无痉挛率达到84.7%，显著高于ACTH组的61.6%，差异具有统计学意义。进一步从总体疗效看，泼尼松和或氨己烯酸组有效率为85.5%，也高于ACTH组的70.6%。这说明从真实世界数据来看，口服方案不仅没有处于劣势，反而在短期临床止痉方面表现得更好。")

    # 13. Recurrence
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "6个月复发率：口服方案复发更少", "在初始缓解患儿中，复发控制优势非常明确")
    add_metric_card(slide, 0.85, 2.0, 2.8, 1.55, "纳入复发分析", "172例", "仅纳入达到初始缓解的患儿")
    add_metric_card(slide, 4.0, 2.0, 2.9, 1.55, "ACTH组复发率", "50.9%", "110例初始缓解患儿中有56例复发", accent=RED, value_color=RED)
    add_metric_card(slide, 7.25, 2.0, 2.9, 1.55, "泼尼松和/或氨己烯酸组", "27.4%", "62例初始缓解患儿中有17例复发", accent=TEAL)
    add_metric_card(slide, 10.35, 2.0, 1.95, 1.55, "统计学", "P = 0.003", "差异显著", accent=CYAN)
    add_card(slide, 0.95, 4.1, 11.35, 2.1, "结论解读", ["口服方案不仅短期止痉更好，而且在6个月内维持缓解方面更稳定。复发率从50.9%下降到27.4%，这使其临床吸引力明显增强。"], accent=ORANGE)
    add_data_note(slide, 0.95, 6.38, "表3.4、图3.4")
    add_page_no(slide, 13)
    record_script("6个月复发情况", ["复发分析只纳入初始缓解患儿", "口服方案组复发率明显更低"], "在达到初始缓解的172例患儿中，我们进一步分析了6个月内复发情况。结果显示，ACTH组复发率为50.9%，而泼尼松和或氨己烯酸组只有27.4%，差异具有明确统计学意义。这个结果非常重要，因为它说明口服方案的优势不只体现在短期止痉，更体现在维持疗效和降低短期复发风险上。")

    # 14. Adverse events
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "不良反应类型相近，感染最常见", "总体不良反应率口服方案更低，但差异未达显著")
    add_column_chart(slide, 0.7, 2.05, 7.1, 4.35, ["易激惹", "胃肠道", "血压升高", "感染", "嗜睡", "肝功异常"], [("ACTH组", [13.3, 15.0, 10.0, 46.7, 0.0, 15.0]), ("泼尼松和/或氨己烯酸组", [13.6, 13.6, 9.1, 40.9, 9.1, 13.6])], max_value=55)
    add_metric_card(slide, 8.2, 2.0, 4.05, 1.35, "总不良反应率", "28.2% vs 40.5%", "口服方案更低，但P = 0.066", accent=ORANGE)
    add_metric_card(slide, 8.2, 3.7, 4.05, 1.35, "最常见不良反应", "感染", "ACTH组46.7%，口服方案组40.9%", accent=RED, value_color=RED)
    add_metric_card(slide, 8.2, 5.4, 4.05, 1.35, "其他观察", "嗜睡仅见于口服方案组", "但总体数量少，临床可管理", accent=CYAN)
    add_data_note(slide, 0.78, 6.45, "表3.5、表3.6、图3.5")
    add_page_no(slide, 14)
    record_script("不良反应比较", ["感染是两组最常见不良反应", "口服方案总体不良反应率更低但未达统计学显著"], "在安全性方面，两组不良反应类型大体相似，最常见的不良反应都是感染。ACTH组感染发生率为46.7%，口服方案组为40.9%。从总体上看，ACTH组不良反应发生率为40.5%，泼尼松和或氨己烯酸组为28.2%，虽然这个差异没有达到统计学显著，但方向上仍提示口服方案可能更安全一些。需要注意的是，口服方案组出现了少量镇静或嗜睡，但总体数量较少，临床可管理。")

    # 15. Economic
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "卫生经济学优势非常明确", "住院天数与住院费用均明显低于ACTH组")
    add_bar_chart(slide, 0.85, 2.25, 4.8, 2.0, ["ACTH组", "泼尼松和/或氨己烯酸组"], [22.67, 11.32], title="首次诊断治疗平均住院天数（天）", color=TEAL, max_value=26)
    add_bar_chart(slide, 0.85, 4.65, 4.8, 2.0, ["ACTH组", "泼尼松和/或氨己烯酸组"], [17769, 7429], title="首次诊断治疗平均住院费用（元）", color=ORANGE, max_value=19000)
    add_metric_card(slide, 6.2, 2.15, 3.05, 1.55, "住院时间降幅", f"{pct_reduction(22.67, 11.32)}%", "从22.67天降至11.32天", accent=TEAL)
    add_metric_card(slide, 9.45, 2.15, 3.05, 1.55, "住院费用降幅", f"{pct_reduction(17769, 7429)}%", "从17769元降至7429元", accent=ORANGE)
    add_card(slide, 6.18, 4.15, 6.32, 2.15, "临床意义", ["在疗效和复发控制不逊于甚至优于ACTH的前提下，口服方案显著减少住院负担和经济成本，对家庭和医疗系统都更友好。"], accent=CYAN)
    add_data_note(slide, 0.88, 6.6, "表3.7、图3.6")
    add_page_no(slide, 15)
    record_script("住院天数与费用", ["住院天数减少约50.1%", "住院费用减少约58.2%", "这是方案推广的重要现实依据"], "卫生经济学是这项研究非常有现实意义的部分。泼尼松和或氨己烯酸组平均住院时间只有11.32天，而ACTH组为22.67天，降幅大约50%。费用方面，ACTH组平均住院费用为17769元，而口服方案组仅为7429元，降幅超过58%。因此，这一方案的优势不只是统计学上的，而是直接体现在患儿家庭的负担减轻和临床可及性的提高上。")

    # 16. Factors
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "疗效影响因素：单因素提示趋势，多因素确认首选药物重要性", "年龄与治疗时机在单因素分析中有一定关联，但独立作用有限")
    add_table(slide, 0.75, 2.1, 5.9, 2.8, [["单因素指标", "结果", "P值"], ["发作年龄 <12月", "有效率 78.6% vs 54.5%", "0.013"], ["治疗启动 <28天", "有效率 80.0% vs 69.3%", "0.090"], ["病因分层", "三组有效率接近", "0.880"], ["性别", "差异不显著", "0.094"]], col_widths=[1.8, 2.7, 1.1], font_size=10)
    add_table(slide, 6.95, 2.1, 5.35, 3.2, [["Logistic变量", "P值", "OR / 结论"], ["性别", "0.366", "无独立意义"], ["发作年龄", "0.084", "无独立意义"], ["治疗时机", "0.244", "无独立意义"], ["首选药物", "0.019", "独立影响因素"]], col_widths=[1.95, 1.0, 1.85], font_size=10)
    add_card(slide, 0.86, 5.35, 11.35, 1.45, "解读", ["多因素分析提示，真正稳定关联疗效的是首选治疗药物；而年龄与治疗时机更可能通过复杂病情背景间接影响结果。"], accent=TEAL)
    add_data_note(slide, 0.84, 6.85, "表3.9 - 表3.14")
    add_page_no(slide, 16)
    record_script("疗效影响因素分析", ["单因素分析显示发作年龄和治疗时机有一定趋势", "多因素Logistic回归显示首选药物为独立影响因素"], "从单因素分析看，发作年龄小于12个月的患儿总体有效率更高，治疗启动早于28天的患儿也呈现出更好的疗效趋势。但进入多因素Logistic回归后，真正保持独立统计学意义的变量是首选治疗药物。这说明治疗方案本身与疗效关系最为稳定，而发作年龄和治疗时机的影响可能受到病因和病情严重程度等其他因素调节。")

    # 17. Neurodevelopment
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "研究结果", "神经发育：药物方案差异不显著，治疗时机更关键", "6个月发育结局更依赖是否早期启动治疗")
    add_metric_card(slide, 0.85, 2.0, 2.65, 1.45, "入院发育分度", "P = 0.105", "两组基线发育水平无显著差异", accent=TEAL)
    add_metric_card(slide, 3.82, 2.0, 2.65, 1.45, "6个月发育分度", "P = 0.135", "不同药物方案差异不显著", accent=ORANGE)
    add_metric_card(slide, 6.79, 2.0, 2.65, 1.45, "治疗时机与入院发育", "P = 0.014", "早期治疗组入院时发育更好", accent=CYAN)
    add_metric_card(slide, 9.76, 2.0, 2.65, 1.45, "治疗时机与6个月发育", "P < 0.001", "早期治疗组短期发育更优", accent=RED, value_color=RED)
    fill_shape(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(4.0), Inches(10.8), Inches(0.05)), LINE)
    add_timeline_item(slide, 0.95, 3.7, "入院时", "不同治疗方案发育分度相近", TEAL)
    add_timeline_item(slide, 4.55, 3.7, "6个月", "药物差异不显著，但早期治疗更优", ORANGE)
    add_timeline_item(slide, 8.15, 3.7, "1年", "治疗时机组间差异不再显著（P = 0.596）", CYAN)
    add_data_note(slide, 0.95, 6.4, "表3.15 - 表3.20、图3.7、图3.8")
    add_page_no(slide, 17)
    record_script("神经发育结果", ["不同药物方案在6个月发育结局上无显著差异", "早期治疗与更好的短期发育结局密切相关", "1年差异减弱"], "在神经发育方面，ACTH组和泼尼松和或氨己烯酸组无论在入院时还是在6个月随访时，发育分级都没有显著差异。这提示单纯药物种类并不是决定短期发育结局的唯一因素。真正更值得关注的是治疗时机。起病后28天内启动治疗的患儿，在入院时和6个月随访时发育结局都更好，提示早期治疗对短期神经发育具有保护作用。不过到1年随访时，这种差异减弱，说明长期预后仍会受到病因、结构异常和持续癫痫活动等多因素影响。")

    # 18. Discussion
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "讨论", "如何理解本研究结果？", "结果与现有文献总体一致，但更突出真实世界的可及性价值")
    add_card(slide, 0.72, 1.95, 3.9, 4.9, "结果为何合理", ["口服激素联合或不联合VGB可同时覆盖激素作用与GABA机制。", "短期止痉更快，可能有助于减少持续异常放电对大脑的损伤。", "联合方案在真实世界中更容易足量、及时启动。"])
    add_card(slide, 4.84, 1.95, 3.9, 4.9, "与文献的一致性", ["本研究结果与ICISS等研究中“联合治疗短期更优”的方向基本一致。", "ACTH并非在所有终点上都占优势，尤其在可及性与成本方面不占优。"], accent=ORANGE)
    add_card(slide, 8.96, 1.95, 3.4, 4.9, "最重要的临床提醒", ["IESS是时间敏感型疾病。", "早诊断、早治疗，比反复犹豫用哪种药更重要。", "病因评估与康复干预要尽早并行。"], accent=CYAN)
    add_page_no(slide, 18)
    record_script("讨论与文献对照", ["解释口服方案为何可能表现更优", "强调本研究与既往联合治疗研究方向一致"], "本研究结果之所以具有临床说服力，一方面是因为口服激素联合或不联合氨己烯酸在机制上可以更快控制发作，另一方面是因为它在真实世界中更容易被及时启动。我们的结果与既往ICISS等研究方向总体一致，也就是联合治疗或强化治疗在短期控制方面往往更优。更重要的是，本研究把真实世界的成本和可及性纳入了比较，这使口服方案的优势更完整地被呈现出来。")

    # 19. Limitations and outlook
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "局限性与展望", "结论可用于指导临床，但仍需更高等级证据验证", "未来重点应放在多中心、前瞻性和长期发育随访")
    add_card(slide, 0.72, 1.95, 4.0, 4.9, "研究局限", ["单中心回顾性设计，存在选择偏倚。", "部分亚组样本量有限，分层分析稳健性受限。", "存在10% - 20%失访，长期结果代表性可能受影响。"], accent=ORANGE)
    add_card(slide, 4.98, 1.95, 4.0, 4.9, "尚未解决的问题", ["最佳治疗时间窗仍缺乏统一界定。", "不同病因亚群的最佳治疗策略可能并不相同。", "长期认知和社交功能结局仍缺乏足够随访证据。"], accent=CYAN)
    add_card(slide, 9.24, 1.95, 3.2, 4.9, "后续方向", ["开展多中心前瞻性研究。", "细化病因分层管理。", "把标准化发育评估纳入长期随访。"], accent=TEAL)
    add_page_no(slide, 19)
    record_script("局限性与展望", ["说明单中心回顾性研究的边界", "指出未来应加强前瞻性和长期随访研究"], "需要说明的是，这项研究仍然存在边界。首先，它是单中心回顾性研究，治疗方案的选择受到当时临床判断影响，难以完全避免选择偏倚。其次，部分亚组样本量较小，而且存在一定失访，可能影响长期结果的稳定性。未来更理想的方向是开展多中心前瞻性研究，同时针对不同病因亚群建立更精细化的治疗路径，并把长期神经发育评估纳入标准随访体系。")

    # 20. Conclusion
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    add_header(slide, "结论", "本研究的五点主要结论", "结论围绕疗效、复发、经济学、影响因素与发育结局展开")
    add_card(slide, 0.72, 1.95, 3.9, 4.9, "结论1-2", ["泼尼松和（或）氨己烯酸方案短期疗效更好，2周无痉挛率和总体有效率均更高。", "该方案6个月复发率更低，维持缓解优势更明显。"], accent=TEAL)
    add_card(slide, 4.84, 1.95, 3.9, 4.9, "结论3-4", ["口服方案住院时间更短、费用更低，卫生经济学优势明确。", "首选药物是疗效独立影响因素，说明治疗方案选择确实重要。"], accent=ORANGE)
    add_card(slide, 8.96, 1.95, 3.4, 4.9, "结论5", ["不同药物方案对短期发育结局总体影响相近。", "治疗时机对6个月神经发育结局更关键，强调早期识别与及时干预。"], accent=CYAN)
    add_page_no(slide, 20)
    record_script("结论", ["用五点结论收束全文", "再次强调口服方案优势和早治疗的重要性"], "最后总结本研究的五点主要结论。第一，泼尼松和或氨己烯酸方案在短期疗效上优于ACTH。第二，这一方案在6个月复发控制方面也更占优势。第三，它具有明显的住院时间和费用优势。第四，首选药物是疗效的独立影响因素。第五，药物方案对短期发育结局的总体影响相近，但治疗时机对6个月神经发育更关键。因此，临床管理应同时重视药物选择和尽早启动治疗。")

    # 21. Thanks
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    fill_shape(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0)), NAVY)
    add_textbox(slide, 1.0, 2.0, 6.4, 0.8, "感谢各位老师聆听", size=28, color=NAVY, bold=True)
    add_textbox(slide, 1.0, 2.9, 6.2, 0.38, "恳请批评指正", size=15, color=TEAL, bold=True)
    add_metric_card(slide, 8.1, 1.85, 4.0, 1.45, "答辩总结", "疗效、复发、成本、时机", "这四个维度共同决定IESS的一线治疗选择", accent=TEAL)
    add_metric_card(slide, 8.1, 3.65, 4.0, 1.45, "实践建议", "早诊断 早治疗", "在可及性允许时，优先考虑综合效益更优的方案", accent=ORANGE)
    record_script("结束页", ["感谢聆听并进入问答环节"], "我的汇报到这里结束。感谢各位老师的聆听，恳请各位老师批评指正。")

    return prs


def build_script_doc():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = DocxInches(0.8)
    section.bottom_margin = DocxInches(0.8)
    section.left_margin = DocxInches(0.9)
    section.right_margin = DocxInches(0.9)

    styles = doc.styles
    for style_name in ["Normal", "Title", "Heading 1", "Heading 2"]:
        style = styles[style_name]
        style.font.name = "Microsoft YaHei"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")

    styles["Normal"].font.size = DocxPt(11)
    styles["Title"].font.size = DocxPt(18)
    styles["Heading 1"].font.size = DocxPt(14)
    styles["Heading 2"].font.size = DocxPt(12)

    title = doc.add_paragraph(style="Title")
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.add_run("IESS论文答辩演讲稿")

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"对应论文：{SOURCE_DOCX.name}\n").bold = True
    meta.add_run(f"对应PPT：{TITLE}\n")
    meta.add_run(f"汇报人：{AUTHOR}    导师：{ADVISOR}    时间：{TIME}")

    intro = doc.add_paragraph()
    intro.add_run("建议时长：").bold = True
    intro.add_run("12 - 15分钟。建议先按PPT页序讲解，再根据答辩老师提问回到相应页补充。")

    for index, item in enumerate(SLIDE_SCRIPTS, start=1):
        heading = doc.add_paragraph(style="Heading 1")
        heading.add_run(f"第{index}页  {item['title']}")

        key_points_heading = doc.add_paragraph(style="Heading 2")
        key_points_heading.add_run("页面要点")
        for point in item["points"]:
            paragraph = doc.add_paragraph()
            paragraph.paragraph_format.left_indent = DocxPt(14)
            paragraph.add_run("• ").bold = True
            paragraph.add_run(point)

        script_heading = doc.add_paragraph(style="Heading 2")
        script_heading.add_run("演讲稿")
        doc.add_paragraph(item["script"])

    return doc


def build_script_text():
    blocks = [
        "IESS论文答辩演讲稿",
        f"对应论文：{SOURCE_DOCX.name}",
        f"对应PPT：{TITLE}",
        f"汇报人：{AUTHOR}    导师：{ADVISOR}    时间：{TIME}",
        "建议时长：12 - 15分钟。",
    ]

    for index, item in enumerate(SLIDE_SCRIPTS, start=1):
        points = "\n".join(f"- {point}" for point in item["points"])
        blocks.append(f"第{index}页  {item['title']}\n页面要点：\n{points}\n演讲稿：\n{item['script']}")

    return "\n\n".join(blocks)


def main():
    _ = Document(SOURCE_DOCX)
    presentation = build_presentation()
    presentation.save(PPT_PATH)

    script_doc = build_script_doc()
    script_doc.save(SCRIPT_DOCX_PATH)
    SCRIPT_TXT_PATH.write_text(build_script_text(), encoding="utf-8")

    check_ppt = Presentation(PPT_PATH)
    check_doc = Document(SCRIPT_DOCX_PATH)

    print(PPT_PATH)
    print(SCRIPT_DOCX_PATH)
    print(SCRIPT_TXT_PATH)
    print(f"slides={len(check_ppt.slides)}")
    print(f"script_paragraphs={len(check_doc.paragraphs)}")


if __name__ == "__main__":
    main()
